# -*- coding: utf-8 -*-
# 工单编号：人工智能CV-AIGC-16-【必选】文旅Agent任务工单-需求分析与功能设计V1.1-20260306
"""
ppt_generator.py - PPT生成模块
功能：使用LLM生成PPT内容大纲 + python-pptx渲染为.pptx文件
支持自定义主题、页数，输出16:9格式专业PPT
"""

import json  # 用于解析LLM输出的JSON大纲
import os  # 用于创建输出目录
from datetime import datetime  # 用于生成带时间戳的文件名
from pathlib import Path  # 路径处理

# python-pptx是可选依赖，导入失败时PPT功能不可用
try:
    from pptx import Presentation  # PPT文档对象
    from pptx.util import Inches, Pt  # 尺寸单位
    from pptx.dml.color import RGBColor  # RGB颜色
    PPTX_AVAILABLE = True  # 标记pptx可用
except ImportError:  # 未安装python-pptx
    PPTX_AVAILABLE = False  # 标记pptx不可用


# ============================================================
# 输出目录配置
# ============================================================
OUTPUT_DIR = Path(__file__).parent / "output"  # PPT输出到研发/output/
OUTPUT_DIR.mkdir(exist_ok=True)  # 确保目录存在


# ============================================================
# 配色方案 - 文旅主题的蓝色+紫色配色
# ============================================================
PRIMARY_COLOR = RGBColor(59, 130, 246)  # 主色：蓝色 #3b82f6
ACCENT_COLOR = RGBColor(139, 92, 246)  # 强调色：紫色 #8b5cf6
DARK_COLOR = RGBColor(30, 41, 59)  # 深色：深蓝灰 #1e293b
WHITE_COLOR = RGBColor(255, 255, 255)  # 白色
LIGHT_TEXT = RGBColor(100, 116, 139)  # 浅色文字


# ============================================================
# 默认PPT内容 - LLM解析失败时的fallback
# ============================================================
FALLBACK_SLIDES = [
    {"title": "文旅智能体方案", "subtitle": "文旅创新智脑 · AI智能体", "bullets": ["项目背景", "核心目标", "实施方案", "预期成果"]},
    {"title": "项目背景", "subtitle": "", "bullets": ["文旅行业数字化转型加速", "游客个性化需求日益增长", "AI技术为文旅带来创新机遇"]},
    {"title": "核心方案", "subtitle": "", "bullets": ["数字人智能导览", "多模态知识检索", "智能运营管理", "创意内容生成"]},
    {"title": "技术架构", "subtitle": "", "bullets": ["LLM大模型驱动", "多模态RAG引擎", "深度图像任务", "Agent智能调度"]},
    {"title": "实施计划", "subtitle": "", "bullets": ["第一期：MVP核心功能", "第二期：深度体验升级", "第三期：生态与智能化"]},
    {"title": "总结与展望", "subtitle": "", "bullets": ["打造文旅AI新体验", "助力文化传承与创新", "谢谢观看！"]},
]  # 6页默认PPT内容


# ============================================================
# 解析LLM输出的JSON大纲
# ============================================================
def parse_outline_json(content: str, slides_count: int, topic: str) -> list:
    """
    从LLM返回的文本中提取PPT大纲JSON
    参数:
        content: LLM返回的原始文本
        slides_count: 期望的页数
        topic: PPT主题
    返回:
        slides列表 [{"title":"...", "subtitle":"...", "bullets":[...]}, ...]
    """
    try:
        # 尝试从文本中提取JSON部分
        json_start = content.find("{")  # 找第一个{
        json_end = content.rfind("}") + 1  # 找最后一个}
        if json_start >= 0 and json_end > json_start:  # 找到了JSON
            content = content[json_start:json_end]  # 截取JSON部分
        outline = json.loads(content)  # 解析JSON
        slides_data = outline.get("slides", [])  # 获取slides数组
        if slides_data:  # 解析成功且非空
            return slides_data  # 返回有效数据
    except (json.JSONDecodeError, KeyError):  # 解析失败
        pass  # 继续使用fallback
    # 返回默认内容
    return FALLBACK_SLIDES  # 兜底方案


# ============================================================
# 生成PPT文件
# ============================================================
def create_pptx(slides_data: list, topic: str) -> str:
    """
    使用python-pptx创建PPT文件
    参数:
        slides_data: 幻灯片数据列表
        topic: PPT主题（用于文件名）
    返回:
        生成的文件名
    """
    prs = Presentation()  # 创建PPT文档
    prs.slide_width = Inches(13.333)  # 16:9宽屏宽度
    prs.slide_height = Inches(7.5)  # 16:9宽屏高度

    for i, slide_data in enumerate(slides_data):  # 遍历每页数据
        slide_layout = prs.slide_layouts[6]  # 选择空白布局
        slide = prs.slides.add_slide(slide_layout)  # 添加新页

        # 设置背景色
        bg = slide.background  # 获取背景对象
        fill = bg.fill  # 获取填充对象
        fill.solid()  # 设置为纯色填充
        if i == 0:  # 首页：深色背景
            fill.fore_color.rgb = DARK_COLOR  # 深蓝灰背景
            text_color = WHITE_COLOR  # 白色文字
        else:  # 内容页：浅色背景
            fill.fore_color.rgb = WHITE_COLOR  # 白色背景
            text_color = DARK_COLOR  # 深色文字

        # 内容页添加左侧彩色装饰条
        if i > 0:  # 非首页
            from pptx.util import Inches as In  # 局部导入简写
            shape = slide.shapes.add_shape(1, In(0), In(0), In(0.15), In(7.5))  # 左侧细条
            shape.fill.solid()  # 纯色填充
            shape.fill.fore_color.rgb = PRIMARY_COLOR if i % 2 == 1 else ACCENT_COLOR  # 奇偶交替颜色
            shape.line.fill.background()  # 无边框

        # 标题位置和样式
        left = Inches(0.8) if i > 0 else Inches(1)  # x坐标
        top = Inches(1.5) if i == 0 else Inches(0.8)  # y坐标
        width = Inches(11)  # 宽度
        height = Inches(1.2)  # 高度

        # 添加标题文本框
        txBox = slide.shapes.add_textbox(left, top, width, height)  # 创建文本框
        tf = txBox.text_frame  # 获取文本框架
        tf.word_wrap = True  # 启用自动换行
        p = tf.paragraphs[0]  # 第一个段落
        p.text = slide_data.get("title", "")  # 设置标题文字
        p.font.size = Pt(40) if i == 0 else Pt(32)  # 首页更大字号
        p.font.bold = True  # 加粗
        p.font.color.rgb = text_color  # 文字颜色

        # 添加副标题（如果有）
        subtitle = slide_data.get("subtitle", "")  # 获取副标题
        if subtitle:  # 有副标题才添加
            txBox2 = slide.shapes.add_textbox(left, top + height + Inches(0.2), width, Inches(0.6))
            tf2 = txBox2.text_frame  # 文本框架
            p2 = tf2.paragraphs[0]  # 段落
            p2.text = subtitle  # 副标题文字
            p2.font.size = Pt(18)  # 字号
            p2.font.color.rgb = LIGHT_TEXT if i > 0 else RGBColor(200, 200, 220)  # 浅色

        # 添加要点列表
        bullets = slide_data.get("bullets", [])  # 获取要点数组
        bullet_top = top + height + Inches(1.0 if subtitle else 0.6)  # 要点起始位置
        txBox3 = slide.shapes.add_textbox(left, bullet_top, Inches(10), Inches(4.5))
        tf3 = txBox3.text_frame  # 文本框架
        tf3.word_wrap = True  # 自动换行

        for j, bullet in enumerate(bullets):  # 遍历要点
            if j == 0:  # 第一个要点
                p = tf3.paragraphs[0]  # 使用首段
            else:  # 后续要点
                p = tf3.add_paragraph()  # 新增段落
            p.text = f"• {bullet}"  # 加圆点前缀
            p.font.size = Pt(18)  # 字号
            p.font.color.rgb = text_color  # 颜色
            p.space_after = Pt(12)  # 段后间距

        # 添加页码
        txBox4 = slide.shapes.add_textbox(Inches(12), Inches(7), Inches(1), Inches(0.3))
        tf4 = txBox4.text_frame  # 文本框架
        p4 = tf4.paragraphs[0]  # 段落
        p4.text = f"{i+1}/{len(slides_data)}"  # 页码格式：1/6
        p4.font.size = Pt(10)  # 小字号
        p4.font.color.rgb = RGBColor(150, 150, 150)  # 灰色

    # 生成文件名并保存
    safe_topic = topic[:20].replace("/", "_").replace("\\", "_")  # 截断并替换非法字符
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")  # 时间戳
    filename = f"文旅智能体_{safe_topic}_{timestamp}.pptx"  # 完整文件名
    filepath = OUTPUT_DIR / filename  # 完整路径
    prs.save(str(filepath))  # 保存PPT文件
    return filename  # 返回文件名供下载
