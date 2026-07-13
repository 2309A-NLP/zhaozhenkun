# -*- coding: utf-8 -*-
# 工单编号：人工智能NLP-Agent数字人项目-17-教育Agent任务工单-教学场景功能分析及智能备课
# 模块：多格式导出服务 - 支持PDF/Word/PPT/Markdown格式的内容导出
# 创建时间：2025年6月
# 作者：Agent智能备课系统开发组

import os  # 文件系统操作
import uuid  # 唯一ID生成
from typing import List, Optional, Dict  # 类型提示
from datetime import datetime  # 时间处理
from config import get_settings, get_data_dirs  # 系统配置


class ExportService:
    """文档导出服务类 - 将教学内容导出为多种格式的文档"""

    def __init__(self):
        """初始化导出服务 - 确保导出目录存在并配置导出工具"""
        self.settings = get_settings()  # 获取系统配置
        get_data_dirs()  # 确保数据目录存在
        self.export_dir = self.settings.EXPORT_DIR  # 导出文件存储目录
        self._check_export_tools()  # 检查可用的导出工具

    def _check_export_tools(self) -> Dict[str, bool]:
        """检查导出工具可用性 - 检测系统是否安装了各格式导出所需的库"""
        tools_available = {"pdf": False, "docx": False, "pptx": False, "markdown": True}  # Markdown总是可用
        try:
            from docx import Document  # 尝试导入python-docx
            tools_available["docx"] = True  # Word导出可用
        except ImportError:  # 库未安装
            print("python-docx未安装，Word导出功能将受限")  # 警告
        try:
            from pptx import Presentation  # 尝试导入python-pptx
            tools_available["pptx"] = True  # PPT导出可用
        except ImportError:  # 库未安装
            print("python-pptx未安装，PPT导出功能将受限")  # 警告
        try:
            from reportlab.lib.pagesizes import A4  # 尝试导入reportlab
            tools_available["pdf"] = True  # PDF导出可用（使用reportlab）
        except ImportError:  # 库未安装
            try:
                import markdown  # 尝试导入markdown
                tools_available["pdf"] = True  # 可通过markdown转PDF
            except ImportError:  # markdown也不可用
                print("PDF导出工具未安装，将使用纯文本降级方案")  # 警告
        self.tools_available = tools_available  # 保存工具可用性状态
        return tools_available  # 返回工具状态

    def export_to_markdown(self, content_list: List[Dict], title: str = "教学内容") -> str:
        """导出为Markdown文件 - 将内容列表拼接为完整的Markdown文档"""
        markdown_lines = [  # Markdown文档头部
            f"# {title}",  # 主标题
            f"",  # 空行
            f"---",  # 分隔线
            f"",  # 空行
            f"> 生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",  # 生成时间
            f"> 工单编号：{self.settings.WORK_ORDER_ID[:40]}",  # 工单编号
            f"> 导出格式：Markdown",  # 格式说明
            f"",  # 空行
            f"---",  # 分隔线
            f"",  # 空行
        ]
        for idx, content in enumerate(content_list):  # 遍历每条内容
            content_type_name = content.get("content_type", "content")  # 内容类型
            # 中文类型名映射
            type_names = {"lesson_plan": "教案", "courseware": "课件",
                          "exercise": "习题", "case_study": "教学案例", "exam_paper": "试卷"}
            cn_type = type_names.get(content_type_name, content_type_name)  # 获取中文名
            markdown_lines.append(f"# {cn_type}：{content.get('title', '未命名')}")  # 章节标题
            markdown_lines.append("")  # 空行
            markdown_lines.append(content.get("raw_content", ""))  # 原始内容
            markdown_lines.append("")  # 空行
            markdown_lines.append("---")  # 分隔线
            markdown_lines.append("")  # 空行
        full_markdown = "\n".join(markdown_lines)  # 拼接完整Markdown
        file_name = f"export_{uuid.uuid4().hex[:8]}_{title}.md".replace(" ", "_")  # 生成文件名
        file_path = os.path.join(self.export_dir, file_name)  # 完整文件路径
        with open(file_path, "w", encoding="utf-8") as f:  # 打开文件写入
            f.write(full_markdown)  # 写入Markdown内容
        return file_path  # 返回文件路径

    def export_to_docx(self, content_list: List[Dict], title: str = "教学内容") -> str:
        """导出为Word文档 - 使用python-docx生成.docx格式文件"""
        if not self.tools_available.get("docx"):  # python-docx不可用
            return self._fallback_export(content_list, title, "docx")  # 降级方案
        from docx import Document  # 导入python-docx
        from docx.shared import Pt, Inches, RGBColor  # 文档样式
        from docx.enum.text import WD_ALIGN_PARAGRAPH  # 文本对齐

        doc = Document()  # 创建Word文档对象
        # 设置文档标题
        title_para = doc.add_heading(title, level=0)  # 添加主标题
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER  # 标题居中
        # 添加元信息
        meta = doc.add_paragraph()  # 创建元信息段落
        meta.add_run(f"生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M')}\n").font.size = Pt(10)  # 时间信息
        meta.add_run(f"工单编号：{self.settings.WORK_ORDER_ID[:50]}").font.size = Pt(10)  # 工单编号
        doc.add_paragraph()  # 空行
        for content in content_list:  # 遍历内容列表
            doc.add_heading(content.get("title", "未命名"), level=1)  # 添加一级标题
            # 处理Markdown内容转为Word段落
            raw_text = content.get("raw_content", "")  # 获取原始Markdown
            paragraphs = raw_text.split("\n\n")  # 按空行分段
            for para_text in paragraphs:  # 遍历每段
                if para_text.strip():  # 段落非空
                    p = doc.add_paragraph()  # 添加段落
                    # 处理Markdown标题格式
                    if para_text.startswith("### "):  # 三级标题
                        p.style = doc.styles["Heading 3"]  # 设置三级标题样式
                        para_text = para_text[4:]  # 去除Markdown标记
                    elif para_text.startswith("## "):  # 二级标题
                        p.style = doc.styles["Heading 2"]  # 设置二级标题样式
                        para_text = para_text[3:]  # 去除Markdown标记
                    p.add_run(para_text.strip()).font.size = Pt(11)  # 添加文本内容
            doc.add_page_break()  # 分页
        file_name = f"export_{uuid.uuid4().hex[:8]}_{title}.docx".replace(" ", "_")  # 生成文件名
        file_path = os.path.join(self.export_dir, file_name)  # 完整路径
        doc.save(file_path)  # 保存Word文档
        return file_path  # 返回文件路径

    def export_to_pptx(self, content_list: List[Dict], title: str = "教学课件") -> str:
        """导出为PPT课件 - 使用python-pptx生成.pptx格式的演示文稿"""
        if not self.tools_available.get("pptx"):  # python-pptx不可用
            return self._fallback_export(content_list, title, "pptx")  # 降级方案
        from pptx import Presentation  # 导入python-pptx
        from pptx.util import Inches, Pt  # PPT尺寸单位
        from pptx.enum.text import PP_ALIGN  # 文本对齐

        prs = Presentation()  # 创建PPT对象
        prs.slide_width = Inches(13.333)  # 设置宽屏16:9
        prs.slide_height = Inches(7.5)  # 设置宽屏高度
        # 封面页 - 安全选择布局（布局0可能不存在于所有模板）
        try:
            slide_layout = prs.slide_layouts[0]  # 尝试使用封面布局
        except (IndexError, KeyError):  # 布局不可用
            slide_layout = prs.slide_layouts[0] if prs.slide_layouts else None  # 降级
        if slide_layout is None:  # 完全没有布局
            slide_layout = prs.slide_layouts[0]  # 强制使用第一个
        slide = prs.slides.add_slide(slide_layout)  # 添加封面幻灯片
        if slide.shapes.title:  # 标题占位符存在
            slide.shapes.title.text = title  # 设置封面标题
        # 安全设置副标题
        subtitle_placeholders = [p for p in slide.placeholders if p.placeholder_format.idx == 1]  # 查找副标题占位符
        if subtitle_placeholders and subtitle_placeholders[0].text_frame:  # 副标题占位符存在
            subtitle_placeholders[0].text = f"智能备课系统自动生成\n{datetime.now().strftime('%Y-%m-%d')}"
        # 内容页
        content_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]  # 安全获取内容布局
        for content in content_list:  # 遍历内容
            raw_text = content.get("raw_content", "")  # 获取原始内容
            sections = raw_text.split("## ")  # 按二级标题分割
            for section in sections[1:]:  # 跳过第一个空段
                slide = prs.slides.add_slide(content_layout)  # 添加幻灯片
                lines = section.strip().split("\n", 1)  # 分割标题和正文
                if slide.shapes.title:  # 标题占位符存在
                    slide.shapes.title.text = lines[0].strip()[:100]  # 设置页标题
                if len(lines) > 1:  # 有正文内容
                    body_placeholders = [p for p in slide.placeholders if p.placeholder_format.idx == 1]  # 查找正文占位符
                    if body_placeholders:  # 正文占位符存在
                        body_placeholders[0].text_frame.text = lines[1].strip()[:500]  # 设置正文（限制长度）
        file_name = f"export_{uuid.uuid4().hex[:8]}_{title}.pptx".replace(" ", "_")  # 生成文件名
        file_path = os.path.join(self.export_dir, file_name)  # 完整路径
        prs.save(file_path)  # 保存PPT文件
        return file_path  # 返回文件路径

    def export_to_pdf(self, content_list: List[Dict], title: str = "教学内容") -> str:
        """导出为PDF文件 - 将Markdown内容转换为PDF格式"""
        # 先生成Markdown
        md_path = self.export_to_markdown(content_list, title)  # 导出为Markdown
        pdf_path = md_path.replace(".md", ".pdf")  # 生成PDF路径
        try:
            import markdown  # 导入markdown库
            with open(md_path, "r", encoding="utf-8") as f:  # 读取Markdown
                md_content = f.read()  # 读取全部内容
            html_content = markdown.markdown(md_content, extensions=["tables", "fenced_code"])  # 转为HTML
            # 使用weasyprint或pdfkit转换HTML到PDF
            try:
                import pdfkit  # 尝试pdfkit
                pdfkit.from_string(html_content, pdf_path)  # HTML转PDF
            except (ImportError, OSError):  # pdfkit不可用
                # 写入HTML文件作为降级方案
                html_path = md_path.replace(".md", ".html")  # HTML路径
                with open(html_path, "w", encoding="utf-8") as f:  # 写入HTML
                    f.write(f"<html><body>{html_content}</body></html>")  # 包裹HTML结构
                return html_path  # 返回HTML路径作为降级
        except ImportError:  # markdown库不可用
            return md_path  # 返回Markdown路径作为降级
        return pdf_path  # 返回PDF路径

    def _fallback_export(self, content_list: List[Dict], title: str, fmt: str) -> str:
        """降级导出 - 当目标格式工具不可用时，导出为Markdown格式作为替代"""
        print(f"{fmt}导出工具不可用，降级为Markdown格式")  # 降级警告
        return self.export_to_markdown(content_list, f"{title}(降级_{fmt})")  # 导出Markdown

    def export_batch(self, content_list: List[Dict], title: str,
                     formats: List[str]) -> Dict[str, str]:
        """批量导出 - 将同一份内容同时导出为多种格式"""
        export_results = {}  # 导出结果字典
        for fmt in formats:  # 遍历目标格式
            print(f"正在导出为 {fmt} 格式...")  # 进度日志
            if fmt == "markdown":  # Markdown格式
                path = self.export_to_markdown(content_list, title)  # 调用Markdown导出
            elif fmt == "docx":  # Word格式
                path = self.export_to_docx(content_list, title)  # 调用Word导出
            elif fmt == "pptx":  # PPT格式
                path = self.export_to_pptx(content_list, title)  # 调用PPT导出
            elif fmt == "pdf":  # PDF格式
                path = self.export_to_pdf(content_list, title)  # 调用PDF导出
            else:  # 不支持的格式
                path = self._fallback_export(content_list, title, fmt)  # 降级导出
            export_results[fmt] = path  # 记录导出路径
            print(f"  ✓ {fmt}: {path}")  # 成功日志
        return export_results  # 返回所有导出结果

    def get_download_info(self, file_path: str) -> Dict:
        """获取下载信息 - 返回文件的下载元数据"""
        if not os.path.exists(file_path):  # 文件不存在
            return {"error": "文件不存在"}  # 返回错误
        file_stat = os.stat(file_path)  # 获取文件状态
        return {"file_path": file_path, "file_name": os.path.basename(file_path),  # 文件名
                "file_size": file_stat.st_size, "file_size_mb": round(file_stat.st_size / 1048576, 2),  # 文件大小
                "format": os.path.splitext(file_path)[1].lstrip("."),  # 文件格式
                "created_at": datetime.fromtimestamp(file_stat.st_ctime).isoformat()}  # 创建时间


# 全局导出服务单例
export_service = ExportService()  # 创建全局唯一的导出服务实例
